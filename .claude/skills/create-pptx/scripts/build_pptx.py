"""
DB Inc. PPTX Builder — DB제목 B/M/L 폰트, 네이비 블루 테마 적용

Usage:
    python build_pptx.py <json_path> <output_path>

JSON 구조:
{
  "title": "프로젝트명",
  "subtitle": "부서명 또는 설명",
  "date": "2025.8",
  "slides": [
    {
      "number": "01",
      "title": "슬라이드 제목",
      "summary": ["요약 1줄", "요약 2줄"],
      "sections": [
        {
          "header": "섹션 헤더 텍스트",
          "type": "table",
          "columns": ["항목", "내용", "비고"],
          "rows": [
            ["라벨1", "값1", "값2"],
            ["라벨2", "값3", "값4"]
          ]
        },
        {
          "header": "섹션 헤더 2",
          "type": "table",
          "columns": ["구분", "설명"],
          "rows": [["A", "B"]]
        }
      ]
    },
    {
      "number": "05",
      "title": "로드맵",
      "summary": ["요약 1줄", "요약 2줄"],
      "sections": [
        {
          "type": "roadmap",
          "phases": [
            {
              "name": "Phase Ⅰ (~3개월)",
              "title": "내부 효율화",
              "description": "핵심 프로세스 선정 및 내부 운영 프로세스 확인",
              "items": ["항목1", "항목2", "항목3"],
              "expected": "테스트 코드 작성 시간 90% 이상 감소"
            }
          ]
        }
      ]
    }
  ]
}
"""

import json
import sys
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Constants ──────────────────────────────────────────────

FONT_BOLD = "DB제목 B"
FONT_MEDIUM = "DB제목 M"
FONT_LIGHT = "DB제목 L"

NAVY_BLUE = RGBColor(0x00, 0x20, 0x60)
ACCENT_BLUE = RGBColor(0x44, 0x72, 0xC4)
DARK_TEXT = RGBColor(0x44, 0x54, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)

# Slide dimensions (widescreen 16:9)
SLIDE_W = 12192000  # EMU
SLIDE_H = 6858000

# Content area
CONTENT_LEFT = Emu(704850)
CONTENT_TOP_SECTION = Emu(1860000)
CONTENT_WIDTH = Emu(10760000)

SCRIPT_DIR = Path(__file__).parent
TEMPLATE_PATH = SCRIPT_DIR.parent / "assets" / "db-template.pptx"


# ── Helper Functions ───────────────────────────────────────

def set_run(run, text, font_name=None, size=None, bold=None, italic=None, color=None):
    """Set run properties immutably-style (returns run for chaining)."""
    run.text = text
    font = run.font
    if font_name:
        font.name = font_name
    if size:
        font.size = Pt(size)
    if bold is not None:
        font.bold = bold
    if italic is not None:
        font.italic = italic
    if color:
        font.color.rgb = color
    return run


def add_run(paragraph, text, font_name=None, size=None, bold=None, color=None):
    """Add a new run to paragraph with specified styling."""
    run = paragraph.add_run()
    return set_run(run, text, font_name, size, bold, color=color)


def add_section_header(slide, text, left, top, width, height=Emu(295275)):
    """Add a rounded rectangle section header bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, text, FONT_BOLD, 12, color=WHITE)

    # Vertical center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return shape


def add_table(slide, columns, rows, left, top, width, col_widths=None):
    """Add a styled table with DB font conventions."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(columns)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Emu(300000))
    table = table_shape.table

    # Auto-calculate column widths if not provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)
    else:
        col_w = int(width / n_cols)
        # First column narrower for labels
        if n_cols >= 2:
            label_w = int(width * 0.18)
            body_w = int((width - label_w) / (n_cols - 1))
            table.columns[0].width = Emu(label_w)
            for i in range(1, n_cols):
                table.columns[i].width = Emu(body_w)
        else:
            table.columns[0].width = Emu(col_w)

    # Header row
    for col_idx, col_text in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, col_text, FONT_MEDIUM, 12, bold=True)
        # Header row fill
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xD6, 0xE4, 0xF0)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]

            # Handle multi-line with pipe separator
            lines = str(cell_text).split("|")
            for line_idx, line in enumerate(lines):
                if line_idx > 0:
                    p = cell.text_frame.add_paragraph()

                if col_idx == 0:
                    # Label column: Bold, Medium font
                    p.alignment = PP_ALIGN.CENTER
                    add_run(p, line.strip(), FONT_MEDIUM, 12, bold=True)
                else:
                    # Body columns: Light font
                    p.alignment = PP_ALIGN.LEFT
                    add_run(p, line.strip(), FONT_LIGHT, 11)

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Set row heights
    header_height = Emu(350000)
    body_height = Emu(450000)
    table.rows[0].height = header_height
    for i in range(1, n_rows):
        table.rows[i].height = body_height

    return table_shape


def add_flow_diagram(slide, steps, top_start):
    """Add horizontal flow diagram with step cards and chevron connectors."""
    n = len(steps)
    if n == 0:
        return int(top_start)

    content_left = int(CONTENT_LEFT)
    content_width = int(CONTENT_WIDTH)

    # Dimensions
    arrow_w = 260000
    n_arrows = n - 1
    total_arrow_space = arrow_w * n_arrows
    box_w = int((content_width - total_arrow_space) / n)

    # Calculate box height from max items
    max_items = max(len(s.get("items", [])) for s in steps)
    badge_h = 280000
    title_h = 320000
    item_h = 230000
    padding = 200000
    box_h = badge_h + title_h + max_items * item_h + padding

    for i, step in enumerate(steps):
        x = content_left + i * (box_w + arrow_w)
        y = int(top_start)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(x), Emu(y), Emu(box_w), Emu(box_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(1.5)

        # Number badge bar at top of card
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(x), Emu(y), Emu(box_w), Emu(badge_h)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_BLUE
        badge.line.fill.background()
        p = badge.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, step.get("label", f"Step {i+1}"), FONT_BOLD, 11, color=WHITE)

        # Title
        title_y = y + badge_h + 60000
        title_box = slide.shapes.add_textbox(
            Emu(x + 60000), Emu(title_y),
            Emu(box_w - 120000), Emu(title_h)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, step.get("title", ""), FONT_BOLD, 14)

        # Detail items
        items = step.get("items", [])
        items_start = title_y + title_h
        for j, item in enumerate(items):
            item_box = slide.shapes.add_textbox(
                Emu(x + 100000), Emu(items_start + j * item_h),
                Emu(box_w - 200000), Emu(item_h)
            )
            tf = item_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            add_run(p, f"• {item}", FONT_LIGHT, 10, color=DARK_TEXT)

        # Chevron arrow between cards
        if i < n - 1:
            arrow_x = x + box_w + 20000
            arrow_y = y + box_h // 2 - 120000
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Emu(arrow_x), Emu(arrow_y),
                Emu(arrow_w - 40000), Emu(240000)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BLUE
            arrow.line.fill.background()

    return int(top_start) + box_h + 200000


def add_roadmap_phases(slide, phases, top_start):
    """Add Phase cards for roadmap slide."""
    n_phases = len(phases)
    if n_phases == 0:
        return

    margin = Emu(100000)
    total_width = int(CONTENT_WIDTH)
    card_width = int((total_width - margin * (n_phases - 1)) / n_phases)

    for i, phase in enumerate(phases):
        left = int(CONTENT_LEFT) + i * (card_width + int(margin))

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left), top_start,
            Emu(card_width), Emu(3800000)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        card.line.width = Pt(0.5)

        # Phase label (top bar)
        label_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left + 200000), Emu(int(top_start) + 200000),
            Emu(card_width - 400000), Emu(276000)
        )
        label_shape.fill.solid()
        label_shape.fill.fore_color.rgb = ACCENT_BLUE
        label_shape.line.fill.background()
        p = label_shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, phase.get("name", f"Phase {i+1}"), FONT_BOLD, 12, color=WHITE)

        # Phase title
        title_top = int(top_start) + 550000
        title_box = slide.shapes.add_textbox(
            Emu(left + 100000), Emu(title_top),
            Emu(card_width - 200000), Emu(340000)
        )
        p = title_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, phase.get("title", ""), FONT_BOLD, 16)

        # Description
        desc_top = title_top + 400000
        desc_box = slide.shapes.add_textbox(
            Emu(left + 100000), Emu(desc_top),
            Emu(card_width - 200000), Emu(280000)
        )
        p = desc_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, phase.get("description", ""), FONT_MEDIUM, 13)

        # Items as mini table
        items = phase.get("items", [])
        if items:
            items_top = desc_top + 350000
            item_height = 260000
            for j, item in enumerate(items):
                item_box = slide.shapes.add_textbox(
                    Emu(left + 200000), Emu(items_top + j * item_height),
                    Emu(card_width - 400000), Emu(item_height)
                )
                p = item_box.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                add_run(p, f"• {item}", FONT_LIGHT, 10)

        # Expected effect box
        expected = phase.get("expected", "")
        if expected:
            effect_top = int(top_start) + 3200000
            effect_label = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(left + int(card_width * 0.25)), Emu(effect_top),
                Emu(int(card_width * 0.5)), Emu(210000)
            )
            effect_label.fill.solid()
            effect_label.fill.fore_color.rgb = NAVY_BLUE
            effect_label.line.fill.background()
            p = effect_label.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            add_run(p, "기대효과", FONT_MEDIUM, 10, bold=True, color=WHITE)

            effect_box = slide.shapes.add_textbox(
                Emu(left + 100000), Emu(effect_top + 260000),
                Emu(card_width - 200000), Emu(400000)
            )
            tf = effect_box.text_frame
            tf.word_wrap = True
            for line in expected.split("|"):
                p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                add_run(p, line.strip(), FONT_MEDIUM, 12, bold=True)


# ── Main Builder ───────────────────────────────────────────

def build_pptx(data, output_path):
    """Build a PPTX file from structured JSON data."""
    # Load template for slide masters/layouts
    if TEMPLATE_PATH.exists():
        prs = Presentation(str(TEMPLATE_PATH))
        # Remove existing slides (keep masters)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            )
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    else:
        prs = Presentation()
        prs.slide_width = Emu(SLIDE_W)
        prs.slide_height = Emu(SLIDE_H)

    # Get layouts
    cover_layout = prs.slide_layouts[0]   # 표지
    content_layout = prs.slide_layouts[1]  # 내용
    eod_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[-1]

    # ── Slide 1: Cover ──
    slide = prs.slides.add_slide(cover_layout)
    # Fill placeholders
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 13:  # Title
            ph.text = data.get("title", "제목을 입력하세요")
        elif idx == 14:  # Subtitle
            ph.text = data.get("subtitle", "")
        elif idx == 15:  # Date
            ph.text = data.get("date", "2025")

    # ── Content Slides ──
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(content_layout)

        # Fill placeholders
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx == 10:  # Title
                ph.text = slide_data.get("title", "")
            elif idx == 11:  # Number
                tf = ph.text_frame
                tf.paragraphs[0].clear()
                add_run(tf.paragraphs[0], slide_data.get("number", "01"), color=NAVY_BLUE)
            elif idx == 12:  # Summary
                tf = ph.text_frame
                tf.paragraphs[0].clear()
                summary = slide_data.get("summary", [])
                for i, line in enumerate(summary):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    add_run(p, line, size=14, bold=True)

        # Sections
        current_top = Emu(1860000)
        sections = slide_data.get("sections", [])

        for section in sections:
            section_type = section.get("type", "table")

            if section_type == "flow":
                # Flow diagram header
                header = section.get("header", "")
                if header:
                    add_section_header(slide, header, CONTENT_LEFT, current_top, CONTENT_WIDTH)
                    current_top = Emu(int(current_top) + 350000)

                steps = section.get("steps", [])
                new_top = add_flow_diagram(slide, steps, current_top)
                current_top = Emu(new_top)

            elif section_type == "roadmap":
                # Roadmap header
                header = section.get("header", "")
                if header:
                    add_section_header(slide, header, CONTENT_LEFT, current_top, CONTENT_WIDTH)
                    current_top = Emu(int(current_top) + 350000)

                phases = section.get("phases", [])
                add_roadmap_phases(slide, phases, current_top)
                current_top = Emu(int(current_top) + 4200000)

            elif section_type == "table":
                # Section header bar
                header = section.get("header", "")
                if header:
                    add_section_header(slide, header, CONTENT_LEFT, current_top, CONTENT_WIDTH)
                    current_top = Emu(int(current_top) + 320000)

                # Table
                columns = section.get("columns", [])
                rows = section.get("rows", [])
                col_widths = section.get("col_widths", None)

                if columns and rows:
                    tbl_shape = add_table(
                        slide, columns, rows,
                        CONTENT_LEFT, current_top, CONTENT_WIDTH,
                        col_widths=col_widths
                    )
                    # Calculate table height for next section positioning
                    n_rows = len(rows) + 1
                    table_height = 350000 + 450000 * len(rows)
                    current_top = Emu(int(current_top) + table_height + 200000)

            elif section_type == "text":
                # Simple text block
                text_box = slide.shapes.add_textbox(
                    CONTENT_LEFT, current_top, CONTENT_WIDTH, Emu(500000)
                )
                tf = text_box.text_frame
                tf.word_wrap = True
                content = section.get("content", "")
                for line in content.split("|"):
                    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
                    add_run(p, line.strip(), FONT_LIGHT, 12)
                current_top = Emu(int(current_top) + 550000)

    # ── EOD Slide ──
    prs.slides.add_slide(eod_layout)

    # Save
    prs.save(output_path)
    print(f"PPTX saved: {output_path}")
    return output_path


# ── CLI Entry Point ────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python build_pptx.py <json_path> <output_path>")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2]

    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    build_pptx(data, output_path)
